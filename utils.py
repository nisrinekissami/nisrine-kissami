"""
utils.py
---------
Fonctions de calcul pour le dashboard de réduction du downtime.
Toutes les fonctions prennent un DataFrame pandas "standardisé" avec au minimum
les colonnes : Date, Machine, Cause, Duree_min
(la colonne Responsable est optionnelle)

Ce module est volontairement séparé de app.py pour que la logique de calcul
soit testable indépendamment de l'interface Streamlit.
"""

import re

import pandas as pd
import numpy as np

# ---------------------------------------------------------------------------
# PALETTE DE MARQUE VERSIGENT — extraite directement du logo (icône orange +
# texte noir). Utilisée PARTOUT (graphiques Plotly de l'app, feuilles Excel
# exportées, présentation PowerPoint) pour une identité visuelle cohérente,
# au lieu de couleurs génériques différentes à chaque endroit.
# Deux formats sont fournis : sans "#" pour openpyxl/python-pptx
# (RGBColor.from_string, PatternFill), avec "#" pour Plotly.
# ---------------------------------------------------------------------------
VERSIGENT_ORANGE = "CD7925"
VERSIGENT_BLACK = "1A1A1A"
VERSIGENT_ORANGE_HEX = "#CD7925"
VERSIGENT_BLACK_HEX = "#1A1A1A"
VERSIGENT_GRAY_HEX = "#8C8C8C"

REQUIRED_COLUMNS = ["Date", "Machine", "Cause", "Duree_min"]

# Mots-clés utilisés pour deviner automatiquement la colonne "Famille" (catégorie
# de machine : Cutting Machine, Kit Seal, Outils, Press...) dans un export RadGrid
# ou tout autre export de pannes. Chaque nouvel export peut avoir des noms de
# colonnes légèrement différents, donc on garde une liste large de candidats.
FAMILY_COLUMN_CANDIDATES = [
    "sub description", "famille", "family", "categorie", "catégorie",
    "category", "plant group", "type de machine", "asset type",
]


def auto_detect_family_column(columns):
    """Devine la colonne 'Famille' (catégorie de machine) dans une liste de colonnes,
    en cherchant les mots-clés les plus probables. Retourne None si rien ne correspond."""
    for candidate in FAMILY_COLUMN_CANDIDATES:
        for col in columns:
            if candidate in str(col).lower():
                return col
    return None


def _detect_header_row(raw: pd.DataFrame, max_scan: int = 10) -> int:
    """
    Cherche, dans les premières lignes du fichier brut, celle qui ressemble le plus
    à une ligne d'en-tête (plusieurs cellules remplies, contenant des mots-clés
    typiques d'un export de pannes comme 'date', 'duration'/'durée', 'asset', etc.).
    """
    keywords = ["date", "duration", "durée", "duree", "asset", "machine", "cause", "description"]
    for i in range(min(max_scan, len(raw))):
        row = raw.iloc[i]
        non_empty = row.notna().sum()
        row_text = " ".join(str(v).lower() for v in row if pd.notna(v))
        if non_empty >= 3 and any(k in row_text for k in keywords):
            return i
    return 0


def load_data(file) -> pd.DataFrame:
    """
    Charge un fichier Excel (.xlsx/.xls) ou CSV et retourne un DataFrame brut.

    Robuste aux différents exports RadGrid : certains fichiers ont l'en-tête
    directement en ligne 1, d'autres ont une ou plusieurs lignes de titre/logo
    au-dessus. Cette fonction détecte automatiquement la vraie ligne d'en-tête
    au lieu de supposer qu'elle est toujours en première position.
    """
    name = getattr(file, "name", str(file))

    if name.lower().endswith(".csv"):
        raw = pd.read_csv(file, header=None)
    elif name.lower().endswith(".xls"):
        raw = pd.read_excel(file, header=None, engine="xlrd")
    else:
        raw = pd.read_excel(file, header=None)

    header_row_idx = _detect_header_row(raw)

    df = raw.iloc[header_row_idx + 1:].copy()
    df.columns = raw.iloc[header_row_idx]

    # Retire les colonnes sans nom (souvent des colonnes vides en bout de fichier)
    df = df.loc[:, df.columns.notna()]
    df.columns = [str(c).strip() for c in df.columns]

    # Dé-doublonne les noms de colonnes identiques (ex: deux colonnes "Trade"
    # dans un export RadGrid). Sans ça, choisir l'une d'elles dans le mapping
    # de l'interface renomme les DEUX colonnes en même temps (pandas renomme
    # par label, pas par position), ce qui corrompt silencieusement les
    # données. On garde le 1er nom tel quel, et on suffixe les suivants.
    seen = {}
    new_cols = []
    for c in df.columns:
        if c not in seen:
            seen[c] = 0
            new_cols.append(c)
        else:
            seen[c] += 1
            new_cols.append(f"{c} ({seen[c]+1})")
    df.columns = new_cols

    df = df.dropna(how="all").reset_index(drop=True)
    return df


def standardize_columns(df: pd.DataFrame, mapping: dict) -> pd.DataFrame:
    """
    Renomme les colonnes du fichier utilisateur vers le schéma standard
    attendu par les fonctions de calcul, à partir d'un mapping choisi
    dans l'interface (ex: {"Date panne": "Date", "Équipement": "Machine", ...}).
    """
    df = df.rename(columns=mapping)
    missing = [c for c in REQUIRED_COLUMNS if c not in df.columns]
    if missing:
        raise ValueError(
            f"Missing columns after mapping: {missing}. "
            f"Please check the column mapping."
        )
    return df


def _fix_ambiguous_dates(dates: pd.Series, max_ecart_jours: int = 45) -> tuple:
    """
    Corrige automatiquement les dates dont le jour et le mois ont été inversés
    à l'export (bug très fréquent avec les exports RadGrid/Excel : une partie
    des lignes garde une date texte "correcte" pendant qu'une autre partie a
    été re-sérialisée par Excel avec la mauvaise convention jour/mois, par ex.
    "01/07/2026" (1 juillet, format JJ/MM) réinterprété comme le 7 janvier
    (format MM/JJ)).

    Principe (indépendant du fichier, ne suppose aucune plage de dates connue
    à l'avance) :
      1. On sépare les dates "sûres" (jour > 12 : aucune ambiguïté possible
         jour/mois) des dates "ambiguës" (jour <= 12 ET mois <= 12 : les deux
         lectures sont valides).
      2. On calcule la date médiane du groupe "sûr" comme référence de la
         période réelle des données.
      3. Toute date ambiguë située à plus de `max_ecart_jours` jours de cette
         référence est retestée en inversant jour et mois. Si l'inversion la
         rapproche nettement de la référence, on l'applique.
      4. S'il n'y a pas assez de dates "sûres" pour établir une référence
         fiable (ex: fichier couvrant un seul mois), on ne touche à rien :
         mieux vaut ne pas corriger que corriger au hasard.

    Retourne (dates_corrigees: pd.Series, nb_corrections: int).
    """
    dates = dates.copy()
    valides = dates.dropna()
    surs = valides[valides.dt.day > 12]

    if len(surs) < 3:
        return dates, 0

    reference = surs.median()
    n_corrected = 0

    for idx, d in valides.items():
        if d.day > 12:
            continue  # non ambigu, on ne touche pas
        ecart_original = abs((d - reference).days)
        if ecart_original <= max_ecart_jours:
            continue  # déjà cohérent avec le reste des données
        try:
            d_inverse = d.replace(month=d.day, day=d.month)
        except ValueError:
            continue  # inversion impossible (ex: mois=13), on laisse tel quel
        if abs((d_inverse - reference).days) < ecart_original:
            dates.at[idx] = d_inverse
            n_corrected += 1

    return dates, n_corrected


# ---------------------------------------------------------------------------
# Codes génériques de ligne de production (SKODA 370, RP 11...) à retirer des
# libellés de machine / panne (colonnes "Machine" / "Task", issues de "Asset
# Description" / "Task Title") : demandé explicitement par l'utilisateur, car
# ces codes n'apportent aucune information utile pour identifier l'équipement
# et alourdissent inutilement TOUS les affichages (Data, Pareto, Detailed
# Graphs, exports Excel/PowerPoint). Le reste du libellé doit rester intact.
# Exemples :
#   "Conveyor SKODA 370,ENGINE 01"        -> "Conveyor ,ENGINE 01"
#   "LP07-ZU01 USW 5 SKODA 370,LEAD PREP" -> "LP07-ZU01 USW 5 ,LEAD PREP"
#   "RP11 USW04 SKODA 370 BODY 01"        -> "USW04 BODY 01"
EXCLUDED_NAME_TOKENS = ["SKODA 370", "SKODA370", "RP 11", "RP11"]


def _token_to_regex(token: str) -> "re.Pattern":
    """Transforme un token comme 'SKODA 370' en motif tolérant aux espaces
    multiples ou absents ('SKODA370', 'SKODA  370', ...), insensible à la casse."""
    parts = token.split()
    pattern = r"\s*".join(re.escape(p) for p in parts)
    return re.compile(pattern, re.IGNORECASE)


_EXCLUDED_NAME_PATTERNS = [_token_to_regex(t) for t in EXCLUDED_NAME_TOKENS]


def clean_machine_name(value):
    """Retire les codes génériques de ligne listés dans EXCLUDED_NAME_TOKENS
    d'un libellé de machine ou de panne, en gardant le reste du texte intact
    (aucune autre réorganisation des espaces/virgules restants). À appliquer à
    toute colonne pouvant être affichée comme libellé (Machine, Task)."""
    if not isinstance(value, str):
        return value
    cleaned = value
    for pattern in _EXCLUDED_NAME_PATTERNS:
        cleaned = pattern.sub("", cleaned)
    # Ne collapse que les espaces multiples résiduels créés par la suppression
    # (ex: "USW04 SKODA 370 BODY 01" -> "USW04  BODY 01" -> "USW04 BODY 01"),
    # sans toucher aux virgules ni à la ponctuation d'origine.
    cleaned = re.sub(r" {2,}", " ", cleaned).strip()
    return cleaned


def clean_data(df: pd.DataFrame) -> pd.DataFrame:
    """Nettoie les types (dates, durées numériques) et retire les lignes invalides.

    Corrige aussi automatiquement les dates ambiguës (jour/mois inversés),
    un bug d'export très fréquent qui, sinon, disperse silencieusement une
    partie des pannes dans de fausses semaines (ex: des pannes de la semaine
    27 qui se retrouvent affichées en semaine 2, 6 ou 10)."""
    df = df.copy()
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    df["Duree_min"] = pd.to_numeric(df["Duree_min"], errors="coerce")
    df["Machine"] = df["Machine"].astype(str).str.strip().apply(clean_machine_name)
    df["Cause"] = df["Cause"].astype(str).str.strip()

    # La colonne "Famille" (catégorie de machine) est optionnelle : on la nettoie
    # seulement si elle a été mappée, sans jamais l'exiger.
    if "Famille" in df.columns:
        df["Famille"] = df["Famille"].astype(str).str.strip()

    # La colonne "Task" (intitulé de panne, ex: Task Title) peut contenir les
    # mêmes codes génériques de ligne que "Machine" -> même nettoyage, pour que
    # TOUS les graphiques Pareto / Detailed Graphs affichent un libellé propre.
    if "Task" in df.columns:
        df["Task"] = df["Task"].astype(str).str.strip().apply(clean_machine_name)

    df["Date"], n_dates_corrigees = _fix_ambiguous_dates(df["Date"])

    n_before = len(df)
    df = df.dropna(subset=["Date", "Duree_min"])
    df = df[df["Duree_min"] >= 0]
    n_after = len(df)

    df.attrs["lignes_supprimees"] = n_before - n_after
    df.attrs["dates_corrigees"] = n_dates_corrigees
    return df


def compute_pareto(df: pd.DataFrame, group_col: str = "Cause") -> pd.DataFrame:
    """Pareto niveau 1 : classement des causes (ou machines) par durée totale cumulée."""
    grouped = df.groupby(group_col)["Duree_min"].sum().sort_values(ascending=False)
    pct = (grouped / grouped.sum()) * 100
    cumul = pct.cumsum()
    result = pd.DataFrame({
        "Duree_totale_min": grouped,
        "Pourcentage_%": pct.round(1),
        "Cumul_%": cumul.round(1),
    })
    result.index.name = group_col
    return result.reset_index()


def compute_pareto_level2(df: pd.DataFrame, top_value: str, level1_col: str = "Cause",
                           level2_col: str = "Machine") -> pd.DataFrame:
    """Pareto niveau 2 : détail par machine (ou autre dimension) pour UNE cause donnée."""
    subset = df[df[level1_col] == top_value]
    if subset.empty:
        return pd.DataFrame(columns=[level2_col, "Duree_totale_min", "Pourcentage_%", "Cumul_%"])
    return compute_pareto(subset, group_col=level2_col)


def compute_frequency_pareto(df: pd.DataFrame, group_col: str = "Machine") -> pd.DataFrame:
    """
    Pareto par FRÉQUENCE : classement par NOMBRE de pannes (et non par durée cumulée),
    comme le tableau « Frequent breakdowns » utilisé en interne chez Versigent.

    Complémentaire du Pareto par durée (compute_pareto) : un équipement peut apparaître
    en tête ici même si chaque panne est courte, dès lors qu'il tombe en panne très
    souvent — ce qui révèle un problème de fiabilité récurrent, différent d'un simple
    gros arrêt ponctuel repéré par le Pareto en durée.

    Retourne les mêmes noms de colonnes que compute_pareto (Pourcentage_%/Cumul_%),
    pour pouvoir réutiliser exactement les mêmes fonctions de graphique/export.
    """
    grouped = df.groupby(group_col).size().sort_values(ascending=False)
    grouped.name = "Nombre_pannes"
    pct = (grouped / grouped.sum()) * 100
    cumul = pct.cumsum()
    result = pd.DataFrame({
        "Nombre_pannes": grouped,
        "Pourcentage_%": pct.round(1),
        "Cumul_%": cumul.round(1),
    })
    result.index.name = group_col
    return result.reset_index()


def compute_paynter(df: pd.DataFrame, group_col: str = "Cause") -> pd.DataFrame:
    """Tableau Paynter : durée totale par cause et par semaine ISO."""
    df = df.copy()
    df["Semaine"] = df["Date"].dt.isocalendar().year.astype(str) + "-S" + \
        df["Date"].dt.isocalendar().week.astype(str).str.zfill(2)
    paynter = df.pivot_table(
        index=group_col, columns="Semaine", values="Duree_min",
        aggfunc="sum", fill_value=0
    )
    # Trie les semaines chronologiquement
    paynter = paynter[sorted(paynter.columns)]
    return paynter


def compute_alerts(df: pd.DataFrame, threshold_min: float, group_col: str = "Machine") -> pd.DataFrame:
    """Retourne les machines (ou causes) dont le cumul de downtime dépasse le seuil."""
    totals = df.groupby(group_col)["Duree_min"].sum().sort_values(ascending=False)
    alerts = totals[totals > threshold_min]
    return alerts.reset_index().rename(columns={"Duree_min": "Duree_totale_min"})


def predict_next_week(paynter: pd.DataFrame, row_name: str) -> dict:
    """
    Simple next-week downtime prediction for a given cause, using linear
    regression over the weekly history. Falls back to a simple average when
    fewer than 3 data points are available.
    """
    if row_name not in paynter.index:
        return {"prediction": None, "methode": "no data"}

    y = paynter.loc[row_name].values.astype(float)
    n = len(y)

    if n < 3:
        pred = float(np.mean(y)) if n > 0 else 0.0
        return {"prediction": round(pred, 1), "methode": "simple average (limited data)"}

    X = np.arange(n).reshape(-1, 1)
    try:
        from sklearn.linear_model import LinearRegression
        model = LinearRegression().fit(X, y)
        pred = model.predict([[n]])[0]
        pred = max(0.0, float(pred))  # a duration cannot be negative
        return {"prediction": round(pred, 1), "methode": "linear regression"}
    except ImportError:
        # Fallback if scikit-learn is unavailable: moving average over the last 4 weeks
        pred = float(np.mean(y[-4:]))
        return {"prediction": round(pred, 1), "methode": "moving average (scikit-learn unavailable)"}


def generate_text_summary(df: pd.DataFrame, pareto1: pd.DataFrame, group_col: str = "Cause") -> str:
    """
    Generates a short English summary, directly reusable in a report or a
    presentation, from the computed results.
    """
    if df.empty or pareto1.empty:
        return "Not enough data to generate a summary."

    nb_events = len(df)
    total_h = round(df["Duree_min"].sum() / 60, 1)
    top_row = pareto1.iloc[0]
    top_name = top_row[group_col]
    top_pct = top_row["Pourcentage_%"]
    top_2_cumul = pareto1.iloc[min(1, len(pareto1) - 1)]["Cumul_%"]

    nb_causes_80 = (pareto1["Cumul_%"] <= 80).sum() + 1
    nb_causes_80 = min(nb_causes_80, len(pareto1))

    resume = (
        f"Over the analyzed period, {nb_events} failure events were recorded, "
        f"representing a total of {total_h} hours of downtime. "
        f"The main cause is \u00ab {top_name} \u00bb, responsible on its own for {top_pct}% "
        f"of total downtime. "
        f"Following the Pareto rule (80/20), {nb_causes_80} cause(s) out of {len(pareto1)} "
        f"already explain {top_2_cumul}% of the problem. "
        f"It is therefore recommended to prioritize corrective actions on \u00ab {top_name} \u00bb "
        f"before addressing secondary causes."
    )
    return resume


def build_action_plan_template(pareto1: pd.DataFrame, group_col: str = "Cause", top_n: int = 5) -> pd.DataFrame:
    """
    Génère un plan d'action pré-rempli avec les causes prioritaires (issues du Pareto),
    dans le même format que le tableau "Action Plan" utilisé en interne chez Versigent
    (Primary Reason Code / Secondary Reason Code + Root Cause / Corrective Action / Effect).
    Les colonnes à remplir manuellement par l'équipe maintenance sont laissées vides.
    """
    if pareto1.empty:
        return pd.DataFrame()

    top = pareto1.head(top_n).copy()
    plan = pd.DataFrame({
        "Priorité": range(1, len(top) + 1),
        "Primary Reason Code": top[group_col],
        "Root Cause Effect (%)": top["Pourcentage_%"],
        "Secondary Reason Code / Root Cause (à compléter)": "",
        "Permanent Corrective Action (à compléter)": "",
        "Responsable (à compléter)": "",
        "Date échéance (à compléter)": "",
        "Statut (à compléter)": "À traiter",
    })
    return plan


def build_action_plan_par_famille(pareto_par_famille_detail: dict, top_n: int = 3) -> pd.DataFrame:
    """
    Construit un plan d'action pré-rempli à partir des équipements les plus critiques
    de CHAQUE famille (top_n par famille : Cutting Machine, Kit Seal, Outils, Press...),
    dans le même esprit que la feuille "ACTION PLAN" utilisée en interne chez Versigent
    (qui liste des équipements précis comme M02, V827G, Crimping BT 752 ST01 — pas des
    codes de cause génériques).

    À utiliser à la place de build_action_plan_template quand une colonne Famille est
    disponible : c'est ce que l'entreprise attend (un plan d'action qui ne mélange pas
    toutes les familles, comme le Pareto lui-même).
    """
    rows = []
    item_nr = 1
    for famille, sous_pareto in pareto_par_famille_detail.items():
        if sous_pareto.empty:
            continue
        top = sous_pareto.head(top_n)
        machine_col = sous_pareto.columns[0]
        for _, r in top.iterrows():
            rows.append({
                "Item Nr": item_nr,
                "Famille": famille,
                "Équipement / Machine": r[machine_col],
                "Durée totale (min)": r["Duree_totale_min"],
                "% du downtime de la famille": r["Pourcentage_%"],
                "Cause racine (à compléter)": "",
                "Action corrective (à compléter)": "",
                "Responsable (à compléter)": "",
                "Date échéance (à compléter)": "",
                "Statut (à compléter)": "À traiter",
            })
            item_nr += 1
    return pd.DataFrame(rows)


def format_excel_sheet(worksheet, n_cols, header_color=VERSIGENT_ORANGE, header_row=1, start_col=1,
                        max_width=45):
    """
    Applique une mise en forme professionnelle à une feuille Excel :
    en-tête coloré et en gras, texte blanc, colonnes ajustées à la largeur du contenu,
    ligne d'en-tête figée. Rend le fichier exporté directement présentable en réunion,
    SANS que l'opérateur ait besoin de retoucher manuellement la largeur des colonnes.
    `header_row` permet de formater une feuille dont l'en-tête du tableau n'est pas en
    ligne 1 (ex: quand une bannière de titre Versigent occupe la ligne 1).
    `start_col` permet de formater un tableau qui ne commence pas en colonne A
    (ex: tableau décalé à droite pour laisser la place à un graphique à gauche).
    `max_width` plafonne la largeur d'une colonne même si son contenu est très long
    (ex: descriptions de panne) — à augmenter pour les feuilles à texte libre plus long.
    """
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter

    header_fill = PatternFill(start_color=header_color, end_color=header_color, fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)

    for col_idx in range(start_col, start_col + n_cols):
        cell = worksheet.cell(row=header_row, column=col_idx)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    worksheet.row_dimensions[header_row].height = 28

    # Ajuste automatiquement la largeur des colonnes selon leur contenu
    for col_idx in range(start_col, start_col + n_cols):
        letter = get_column_letter(col_idx)
        max_len = 10
        for row in worksheet.iter_rows(min_col=col_idx, max_col=col_idx):
            for cell in row:
                if cell.value is not None:
                    max_len = max(max_len, len(str(cell.value)))
        worksheet.column_dimensions[letter].width = min(max_len + 3, max_width)

    worksheet.freeze_panes = f"A{header_row + 1}"


def add_seuil80_column(worksheet, n_rows, table_col_count, start_row=1, seuil_value=80):
    """
    Ajoute une colonne "Seuil_80" (valeur constante 80) juste à droite d'un tableau de
    Pareto déjà écrit dans la feuille (colonnes 1..table_col_count, lignes
    start_row+1..start_row+n_rows pour les données, start_row pour l'en-tête).
    Sert à tracer la ligne de seuil 80% dans le graphique Pareto associé
    (add_pareto_excel_chart). Retourne l'index de colonne de la nouvelle colonne.
    """
    seuil_col = table_col_count + 1
    worksheet.cell(row=start_row, column=seuil_col, value="Seuil_80")
    for i in range(1, n_rows + 1):
        worksheet.cell(row=start_row + i, column=seuil_col, value=seuil_value)
    return seuil_col


def add_pareto_excel_chart(worksheet, header_row, n_rows, cat_col, val_col, cumul_col,
                            title, anchor, seuil_col=None, value_axis_title="Durée totale (min)",
                            width=17, height=9):
    """
    Insère un vrai graphique Pareto Excel natif (barres = valeur + courbe de cumul %
    sur axe secondaire, avec seuil 80% en pointillé si `seuil_col` est fourni) dans la
    feuille `worksheet`, à partir d'un tableau déjà écrit (en-tête à `header_row`,
    `n_rows` lignes de données juste en dessous). `cat_col`/`val_col`/`cumul_col`/
    `seuil_col` sont des index de colonnes 1-based (1=A, 2=B...), typiquement ceux
    renvoyés par compute_pareto / compute_frequency_pareto : 1=catégorie,
    2=valeur, 4=Cumul_%, et l'éventuelle colonne ajoutée par add_seuil80_column.

    C'est ce graphique natif (et non une simple image collée) qui rend le fichier
    Excel exporté fidèle au format attendu par l'entreprise : un vrai Pareto
    (barres + courbe cumulée + seuil 80%) que l'utilisateur peut rouvrir et
    modifier directement dans Excel.
    """
    from openpyxl.chart import BarChart, LineChart, Reference

    last_row = header_row + n_rows

    bar = BarChart()
    bar.type = "col"
    bar.title = title
    bar.y_axis.title = value_axis_title
    bar.y_axis.majorGridlines = None
    bar.gapWidth = 40
    bar.x_axis.delete = False

    val_ref = Reference(worksheet, min_col=val_col, min_row=header_row, max_row=last_row)
    cat_ref = Reference(worksheet, min_col=cat_col, min_row=header_row + 1, max_row=last_row)
    bar.add_data(val_ref, titles_from_data=True)
    bar.set_categories(cat_ref)
    bar.series[0].graphicalProperties.solidFill = VERSIGENT_ORANGE
    bar.series[0].graphicalProperties.line.noFill = True

    line = LineChart()
    cumul_ref = Reference(worksheet, min_col=cumul_col, min_row=header_row, max_row=last_row)
    line.add_data(cumul_ref, titles_from_data=True)
    line.y_axis.axId = 200
    line.y_axis.title = "Cumul %"
    line.y_axis.scaling.min = 0
    line.y_axis.scaling.max = 100
    line.y_axis.crosses = "max"
    line.y_axis.majorGridlines = None

    if seuil_col is not None:
        seuil_ref = Reference(worksheet, min_col=seuil_col, min_row=header_row, max_row=last_row)
        line.add_data(seuil_ref, titles_from_data=True)

    cumul_series = line.series[0]
    cumul_series.marker.symbol = "circle"
    cumul_series.graphicalProperties.line.solidFill = VERSIGENT_BLACK
    cumul_series.graphicalProperties.line.width = 20000
    cumul_series.smooth = False

    if seuil_col is not None:
        seuil_series = line.series[1]
        seuil_series.marker.symbol = "none"
        seuil_series.graphicalProperties.line.solidFill = "999999"
        seuil_series.graphicalProperties.line.dashStyle = "dash"
        seuil_series.smooth = False

    bar += line
    bar.width = width
    bar.height = height
    worksheet.add_chart(bar, anchor)


def add_repartition_chart_excel(worksheet, header_row, n_rows, cat_col, val_col, title,
                                 anchor, width=11, height=8, donut=True):
    """
    Insère un graphique de répartition (donut par défaut, comme le donut "Répartition
    du downtime par famille" du fichier Cutting Analysis) à partir d'un tableau déjà
    écrit dans la feuille, en-tête à `header_row`, `n_rows` lignes de données.
    Les parts utilisent la palette de marque Versigent (nuances d'orange et de noir)
    au lieu des couleurs arc-en-ciel par défaut d'Excel.
    """
    from openpyxl.chart import PieChart, DoughnutChart, Reference
    from openpyxl.chart.marker import DataPoint
    from openpyxl.chart.shapes import GraphicalProperties

    chart = DoughnutChart() if donut else PieChart()
    chart.title = title
    last_row = header_row + n_rows
    val_ref = Reference(worksheet, min_col=val_col, min_row=header_row, max_row=last_row)
    cat_ref = Reference(worksheet, min_col=cat_col, min_row=header_row + 1, max_row=last_row)
    chart.add_data(val_ref, titles_from_data=True)
    chart.set_categories(cat_ref)

    palette = [VERSIGENT_ORANGE, VERSIGENT_BLACK, "E8A867", "595959", "B5651D", "BFBFBF"]
    chart.series[0].data_points = [
        DataPoint(idx=i, spPr=GraphicalProperties(solidFill=palette[i % len(palette)]))
        for i in range(n_rows)
    ]

    chart.width = width
    chart.height = height
    worksheet.add_chart(chart, anchor)


def summary_kpis(df: pd.DataFrame) -> dict:
    """Indicateurs clés affichés en haut du dashboard."""
    nb_evenements = len(df)
    duree_totale_min = round(df["Duree_min"].sum(), 1) if not df.empty else 0.0
    return {
        "nb_evenements": nb_evenements,
        "duree_totale_min": duree_totale_min,
        "duree_totale_h": round(df["Duree_min"].sum() / 60, 1) if not df.empty else 0.0,
        # Average = total downtime / number of failures (sum divided by count),
        # shown in the interface next to "Number of failures" and "Total downtime".
        "duree_moyenne_min": round(duree_totale_min / nb_evenements, 1) if nb_evenements else 0.0,
        "machine_top": df.groupby("Machine")["Duree_min"].sum().idxmax() if not df.empty else "N/A",
        "cause_top": df.groupby("Cause")["Duree_min"].sum().idxmax() if not df.empty else "N/A",
        "periode_debut": df["Date"].min(),
        "periode_fin": df["Date"].max(),
    }


# ---------------------------------------------------------------------------
# EXPORT POWERPOINT — reprend TOUS les graphiques (images PNG déjà générées par
# l'app avec Plotly/kaleido, donc mêmes couleurs et mêmes titres) et tableaux du
# dashboard dans une présentation prête à partager, sans avoir à ouvrir Excel.
# ---------------------------------------------------------------------------

THEME_COLOR = VERSIGENT_ORANGE   # orange du logo Versigent — en-têtes Excel, titres PPT
ACCENT_COLOR = VERSIGENT_BLACK   # noir du logo Versigent — plan d'action, contraste
LIGHT_BG = "FBF0E4"              # nuance d'orange très claire, pour les cartes KPI


def _pptx_add_table(slide, df, left, top, width, max_rows=12, header_color=THEME_COLOR):
    """Ajoute un tableau PowerPoint natif (pas une image) à partir d'un DataFrame,
    avec le même en-tête coloré/gras que les feuilles Excel exportées."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    df_show = df.head(max_rows)
    n_rows, n_cols = df_show.shape[0] + 1, df_show.shape[1]
    row_h = Inches(0.35)
    height = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    header_rgb = RGBColor.from_string(header_color)
    for j, col_name in enumerate(df_show.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i in range(df_show.shape[0]):
        for j in range(df_show.shape[1]):
            cell = table.cell(i + 1, j)
            val = df_show.iloc[i, j]
            cell.text = "" if pd.isna(val) else str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)

    bottom = top + height
    if len(df) > max_rows:
        note = slide.shapes.add_textbox(left, bottom + Inches(0.03), width, Inches(0.3))
        note.text_frame.text = (
            f"… {len(df) - max_rows} ligne(s) supplémentaire(s) — voir l'export Excel "
            f"pour le détail complet."
        )
        note.text_frame.paragraphs[0].font.size = Pt(9)
        note.text_frame.paragraphs[0].font.italic = True
        note.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x70, 0x70, 0x70)
        bottom += Inches(0.3)
    return bottom


def build_pptx_report(sections: list, title: str = None, subtitle: str = None,
                       closing_text: str = None, logo_bytes: bytes = None) -> "BytesIO":
    """
    Construit une présentation PowerPoint (.pptx) prête à télécharger, reprenant
    TOUS les graphiques et tableaux déjà calculés/affichés par le dashboard —
    mêmes titres, mêmes légendes — avec l'identité visuelle Versigent (orange
    #CD7925 / noir #1A1A1A du logo, fond blanc, logo sur la slide de titre et
    en petit sur chaque slide de contenu) pour que l'équipe puisse partager les
    résultats en réunion sans avoir à ouvrir Excel.

    `sections` est une liste de dicts décrivant chaque slide de contenu, avec les
    clés optionnelles suivantes :
      - "title"   : titre de la slide (str) — reprend en général le même texte
                    que le st.header()/st.subheader() correspondant dans l'app.
      - "caption" : texte d'explication sous le titre (str) — reprend en général
                    le même texte que le st.caption() correspondant.
      - "kpis"    : liste de tuples (label, valeur) affichés en grandes cartes
                    (utilisé pour la slide de synthèse).
      - "image"   : bytes PNG d'un graphique déjà généré par l'app (fig.to_image).
      - "image2"  : bytes PNG d'un second graphique, affiché à côté du premier
                    (ex : Pareto + camembert de répartition, comme dans l'app).
      - "table"   : DataFrame affiché en tableau natif PowerPoint (tronqué aux
                    premières lignes si besoin, avec une note pour indiquer que
                    le détail complet est dans l'export Excel).
      - "table_header_color" : couleur d'en-tête du tableau pour cette slide
                    (ex : noir #1A1A1A pour la slide Plan d'action).

    `logo_bytes` : bytes PNG du logo Versigent (optionnel). S'il est fourni, il
    est affiché en grand sur la slide de titre/clôture et en petit (coin haut
    droit) sur chaque slide de contenu.

    Retourne un BytesIO prêt à passer à st.download_button.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from io import BytesIO

    THEME = RGBColor.from_string(THEME_COLOR)
    ACCENT = RGBColor.from_string(ACCENT_COLOR)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    LIGHT = RGBColor.from_string(LIGHT_BG)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def _white_bg(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _accent_bar(slide, height=Inches(0.12)):
        """Fine bande orange en bas de slide — touche de marque discrète."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - height, prs.slide_width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _add_logo(slide, width, left, top):
        if not logo_bytes:
            return None
        return slide.shapes.add_picture(BytesIO(logo_bytes), left, top, width=width)

    # --- Slide de titre : logo Versigent + titre + sous-titre, fond blanc ---
    # Sautée si ni titre ni logo ne sont fournis.
    if title or logo_bytes:
        slide = prs.slides.add_slide(blank_layout)
        _white_bg(slide)
        _accent_bar(slide)
        y = Inches(2.3)
        if logo_bytes:
            logo_w = Inches(5.2)
            logo_left = (prs.slide_width - logo_w) / 2
            _add_logo(slide, logo_w, Emu(int(logo_left)), y)
            y = y + Inches(1.15)
        if title:
            tb = slide.shapes.add_textbox(Inches(0.9), y, Inches(11.5), Inches(0.9))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = title
            tf.paragraphs[0].font.size = Pt(30)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = THEME
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            y = y + Inches(0.85)
        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.9), y, Inches(11.5), Inches(0.6))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            tf2.text = subtitle
            tf2.paragraphs[0].font.size = Pt(15)
            tf2.paragraphs[0].font.color.rgb = RGBColor(0x60, 0x60, 0x60)
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slides de contenu ---
    for section in sections:
        slide = prs.slides.add_slide(blank_layout)
        _white_bg(slide)

        title_w = Inches(10.6) if logo_bytes else Inches(12.3)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), title_w, Inches(0.65))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        ttf.text = section.get("title", "")
        ttf.paragraphs[0].font.size = Pt(24)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = THEME

        if logo_bytes:
            _add_logo(slide, Inches(1.3), Inches(11.5), Inches(0.32))

        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95),
                                       Inches(12.3), Pt(1.5))
        rule.fill.solid()
        rule.fill.fore_color.rgb = THEME
        rule.line.fill.background()
        rule.shadow.inherit = False

        y = Inches(1.15)
        if section.get("caption"):
            cap_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), Inches(0.8))
            ctf = cap_box.text_frame
            ctf.word_wrap = True
            ctf.text = section["caption"]
            ctf.paragraphs[0].font.size = Pt(12)
            ctf.paragraphs[0].font.italic = True
            ctf.paragraphs[0].font.color.rgb = RGBColor(0x60, 0x60, 0x60)
            y = Inches(2.0)

        if section.get("kpis"):
            kpis_list = section["kpis"]
            n = len(kpis_list)
            gap = Inches(0.25)
            total_w = Inches(12.3)
            box_w = Emu(int((total_w - gap * (n - 1)) / n))
            x = Inches(0.5)
            for label, value in kpis_list:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, Inches(1.5))
                card.fill.solid()
                card.fill.fore_color.rgb = LIGHT
                card.line.color.rgb = THEME
                card.line.width = Pt(1)
                card.shadow.inherit = False
                vtf = card.text_frame
                vtf.word_wrap = True
                vtf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
                vtf.text = str(value)
                vtf.paragraphs[0].font.size = Pt(24)
                vtf.paragraphs[0].font.bold = True
                vtf.paragraphs[0].font.color.rgb = THEME
                vtf.paragraphs[0].alignment = PP_ALIGN.CENTER
                p2 = vtf.add_paragraph()
                p2.text = label
                p2.font.size = Pt(11)
                p2.font.color.rgb = DARK
                p2.alignment = PP_ALIGN.CENTER
                x = Emu(int(x) + int(box_w) + int(gap))
            y = y + Inches(1.75)

        image_bytes = section.get("image")
        image2_bytes = section.get("image2")
        table_df = section.get("table")
        header_color = section.get("table_header_color", THEME_COLOR)

        if image_bytes and image2_bytes:
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.4), y, width=Inches(6.1))
            slide.shapes.add_picture(BytesIO(image2_bytes), Inches(6.75), y, width=Inches(6.1))
        elif image_bytes and table_df is not None:
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.4), y, width=Inches(7.4))
            _pptx_add_table(slide, table_df, Inches(8.0), y, Inches(4.9), max_rows=10,
                             header_color=header_color)
        elif image_bytes:
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(1.15), y, width=Inches(11.0))
        elif table_df is not None:
            bottom = _pptx_add_table(slide, table_df, Inches(0.7), y, Inches(11.9), max_rows=14,
                                      header_color=header_color)
            y = bottom + Inches(0.15)

        if section.get("body_text"):
            body_box = slide.shapes.add_textbox(Inches(0.7), y, Inches(11.9), Inches(7.3) - y)
            btf = body_box.text_frame
            btf.word_wrap = True
            btf.text = section["body_text"]
            btf.paragraphs[0].font.size = Pt(15)
            btf.paragraphs[0].font.color.rgb = DARK
            btf.paragraphs[0].line_spacing = 1.25

    # --- Slide de clôture : logo + texte, fond blanc, même bande orange ---
    if closing_text:
        slide = prs.slides.add_slide(blank_layout)
        _white_bg(slide)
        _accent_bar(slide)
        y = Inches(2.6)
        if logo_bytes:
            logo_w = Inches(3.4)
            logo_left = (prs.slide_width - logo_w) / 2
            _add_logo(slide, logo_w, Emu(int(logo_left)), y)
            y = y + Inches(0.9)
        tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(11.0), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = closing_text
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# EXTENSION V2 — réplique la structure du classeur Excel Versigent
# ("Analysis", "Pareto - Top 10", "Action Plan", "Detailed Graphs", "Data")
# et ajoute la sélection en cascade Phase (Assemblage / Cutting) -> Famille,
# au lieu d'un simple choix de colonne Famille comme dans la V1.
# =============================================================================

# Mots-clés utilisés pour classer chaque panne dans une PHASE (Assemblage / Cutting)
# à partir de la colonne hiérarchique "Position" (ex: "Building,Manufacturing
# Assembly_FA,SKODA 370,BODY 01" ou "Building,Cutting_CT,...", "Die Center_DC,...").
# Cette colonne est différente de "Asset Description"/"Sub Description" : c'est le
# chemin complet de l'équipement dans l'arborescence usine.
PHASE_KEYWORDS = {
    "Assembly": ["assembly", "assemblage"],
    "Cutting": ["cutting", "die center", "die centre", "lead prep", "press"],
}

POSITION_COLUMN_CANDIDATES = ["position"]


def auto_detect_position_column(columns):
    """Devine la colonne hiérarchique 'Position' (ex: 'Building,Manufacturing
    Assembly_FA,SKODA 370,BODY 01'), utilisée pour déduire la Phase (Assemblage/
    Cutting). Ne doit PAS matcher 'Asset / Position', 'Position 3', 'Position 4'."""
    for col in columns:
        c = str(col).strip().lower()
        if c == "position":
            return col
    return None


def detect_phase(position_value) -> str:
    """Déduit la Phase (Assembly / Cutting / Other) à partir d'une valeur de la
    colonne hiérarchique Position. Robuste : si aucun mot-clé ne correspond,
    retourne 'Other' plutôt que de deviner au hasard."""
    if not isinstance(position_value, str) or not position_value.strip():
        return "Other"
    parts = [p.strip() for p in position_value.split(",")]
    token = parts[1] if len(parts) > 1 else position_value
    token_low = token.lower()
    for phase, keywords in PHASE_KEYWORDS.items():
        if any(k in token_low for k in keywords):
            return phase
    return "Other"


def add_phase_column(df: pd.DataFrame, raw_position_series: pd.Series = None) -> pd.DataFrame:
    """Ajoute une colonne 'Phase' (Assembly / Cutting / Other) au DataFrame déjà
    standardisé, à partir de la colonne brute 'Position' (hiérarchie complète).
    Si `raw_position_series` n'est pas fourni (fichier sans cette colonne), toutes
    les lignes sont classées 'Other' et la sélection Phase/Famille se comporte
    alors comme une simple sélection de Famille (comportement V1 conservé)."""
    df = df.copy()
    if raw_position_series is not None:
        df["Phase"] = raw_position_series.reindex(df.index).map(detect_phase)
        df["Phase"] = df["Phase"].fillna("Other")
    else:
        df["Phase"] = "Other"
    return df


def compute_pareto_with_total(df: pd.DataFrame, group_col: str, value_col: str = "Duree_min",
                               agg: str = "sum") -> pd.DataFrame:
    """Comme un Pareto simple (classement décroissant), mais avec une ligne
    'Total Result' ajoutée à la fin — exactement le format utilisé dans les
    graphiques Excel Versigent ('Pareto - Top 10', graphiques par famille).
    agg='sum' pour la durée totale, agg='count' pour le nombre de pannes (NB Visits).
    """
    if df.empty:
        return pd.DataFrame(columns=[group_col, "Valeur"])
    if agg == "count":
        grouped = df.groupby(group_col).size().sort_values(ascending=False)
    else:
        grouped = df.groupby(group_col)[value_col].sum().sort_values(ascending=False)
    result = grouped.reset_index()
    result.columns = [group_col, "Valeur"]
    total_row = pd.DataFrame({group_col: ["Total Result"], "Valeur": [grouped.sum()]})
    return pd.concat([result, total_row], ignore_index=True)


def compute_nb_visits(df: pd.DataFrame, group_col: str = "Machine") -> pd.DataFrame:
    """Nombre d'interventions (visites) par équipement, avec ligne 'Total Result' —
    réplique le graphique 'NB Visits' de la feuille Analysis Versigent."""
    return compute_pareto_with_total(df, group_col=group_col, agg="count")


def compute_family_weekly_charts(df: pd.DataFrame, family_col: str = "Famille",
                                  machine_col: str = "Machine", max_families: int = 8) -> dict:
    """Reproduit les petits graphiques par famille de la feuille 'F. [Phase] Weekly
    Analysis' Excel : pour chaque famille (ROB, CM, Ultrasonic...), un Pareto par
    équipement (avec ligne Total Result). Les familles avec très peu de volume sont
    regroupées ensemble sous 'Autres' seulement à l'affichage (pas ici), pour rester
    simple et réutilisable telle quelle dans le graphique global de répartition."""
    result = {}
    if family_col not in df.columns or df.empty:
        return result
    for famille, sub in df.groupby(family_col):
        if not str(famille).strip() or str(famille).lower() == "nan":
            continue
        result[famille] = compute_pareto_with_total(sub, group_col=machine_col)
    return result


def compute_detailed_breakdown(df: pd.DataFrame, family_col: str = "Famille",
                                machine_col: str = "Machine", task_col: str = "Task",
                                value_col: str = "Duree_min") -> dict:
    """Réplique la feuille 'Detailed Graphs' : pour chaque famille, le détail
    équipement -> intitulé de panne (Task Title), avec la durée de chacun. Retourne
    un dict {famille: DataFrame[Machine, Task, Duree_min]} trié par équipement
    (du plus impactant au moins impactant) puis par durée décroissante à l'intérieur
    de chaque équipement — pour un graphique à barres horizontales à deux niveaux
    (bracket Machine > Task), comme dans le classeur Excel."""
    result = {}
    if family_col not in df.columns or df.empty:
        return result
    has_task = task_col in df.columns
    for famille, sub in df.groupby(family_col):
        if not str(famille).strip() or str(famille).lower() == "nan":
            continue
        if has_task:
            detail = sub.groupby([machine_col, task_col])[value_col].sum().reset_index()
            detail.columns = ["Machine", "Task", "Duree_min"]
        else:
            detail = sub.groupby(machine_col)[value_col].sum().reset_index()
            detail.columns = ["Machine", "Duree_min"]
            detail["Task"] = detail["Machine"]
        machine_totals = detail.groupby("Machine")["Duree_min"].sum().sort_values(ascending=False)
        detail["_machine_order"] = detail["Machine"].map(machine_totals)
        detail = detail.sort_values(["_machine_order", "Duree_min"], ascending=[False, False])
        detail = detail.drop(columns="_machine_order")
        result[famille] = detail
    return result


VERSIGENT_ACTION_PLAN_COLUMNS = [
    "Item Nr", "Problem Description", "Occurrence Date", "Cause", "Action Plan",
    "Action Level", "Resp", "Status",
    "Effectivity Validation Method", "Validation Date", "Resp ", "Status ",
]


def build_versigent_action_plan(df: pd.DataFrame, machine_col: str = "Machine",
                                 cause_col: str = "Cause", task_col: str = "Task",
                                 date_col: str = "Date", value_col: str = "Duree_min",
                                 top_n: int = 6) -> pd.DataFrame:
    """Génère un plan d'action dans le MÊME format exact que la feuille Excel
    'F. [Phase] Action Plan' Versigent (colonnes Item Nr / Problem Description /
    Occurrence Date / Cause / Action Plan / Action Level / Date of Implement / Resp /
    Status / Effectivity Validation Method / Validation Date / Resp / Status).

    Une ligne par équipement le plus impactant (top_n par durée cumulée) ; les
    colonnes que l'équipe maintenance doit remplir manuellement (Action Plan,
    Action Level, Resp, Status...) sont laissées vides."""
    if df.empty:
        return pd.DataFrame(columns=VERSIGENT_ACTION_PLAN_COLUMNS)

    totals = df.groupby(machine_col)[value_col].sum().sort_values(ascending=False)
    top_machines = totals.head(top_n).index.tolist()

    rows = []
    for machine in top_machines:
        sub = df[df[machine_col] == machine]
        problem_desc = ""
        if task_col in sub.columns and not sub[task_col].dropna().empty:
            problem_desc = sub[task_col].value_counts().idxmax()
        last_date = sub[date_col].max() if date_col in sub.columns else None
        last_date_str = last_date.strftime("%Y-%m-%d") if pd.notna(last_date) else ""
        rows.append({
            "Item Nr": machine,
            "Problem Description": problem_desc,
            "Occurrence Date": last_date_str,
            "Cause": "",  # left blank for the operator to fill in
            "Action Plan": "",
            "Action Level": "",
            "Resp": "",
            "Status": "",
            "Effectivity Validation Method": "Visual Control",
            "Validation Date": "",
            "Resp ": "",
            "Status ": "",
        })
    return pd.DataFrame(rows, columns=VERSIGENT_ACTION_PLAN_COLUMNS)


# =============================================================================
# EXTENSION V3 — regroupement Équipement (CM+Convoyeur / ROB / Ultrasonic /
# Autres), exclusion semaine 24, Plan d'action au format exact Versigent
# (WKxx-j, Action Level=3, Resp=Ingénieur formaté, Ahmed BOUBEGHLI, Done),
# et graphique de répartition avec étiquettes de pourcentage à l'EXTÉRIEUR.
# =============================================================================

# Règle de regroupement demandée : quel que soit le fichier, on ne travaille
# plus avec la colonne Famille brute (ex: 20 catégories différentes) mais avec
# 4 grands groupes d'équipements : CM (qui englobe aussi les convoyeurs),
# ROB, Ultrasonic, et Autres (tout le reste : Torque, DS20, UCAB, Kit Seal,
# Press, Outils, Cutting Machine...). "CM" reste affiché sous le nom "CM"
# même s'il contient aussi les convoyeurs.
EQUIPEMENT_BUCKETS = {
    "CM": ["cm", "conveyor", "convoyeur", "convoyer"],
    "ROB": ["rob"],
    "Ultrasonic": ["ultrasonic", "ultrason", "us"],
}


def map_equipement_bucket(famille_value) -> str:
    """Classe une valeur brute de Famille/Sub Description dans l'un des 4 groupes
    d'équipements demandés : CM, ROB, Ultrasonic, ou Autres."""
    if not isinstance(famille_value, str) or not famille_value.strip():
        return "Others"
    v = famille_value.strip().lower()
    for bucket, keywords in EQUIPEMENT_BUCKETS.items():
        if v in keywords or any(v.startswith(k) for k in keywords):
            return bucket
    return "Others"


def add_equipement_column(df: pd.DataFrame, famille_col: str = "Famille") -> pd.DataFrame:
    """Ajoute la colonne 'Equipement' (CM / ROB / Ultrasonic / Autres) utilisée
    PARTOUT dans l'application à la place de la Famille brute."""
    df = df.copy()
    if famille_col in df.columns:
        df["Equipement"] = df[famille_col].apply(map_equipement_bucket)
    else:
        df["Equipement"] = "Others"
    return df


def exclude_week(df: pd.DataFrame, week_number: int = 24, date_col: str = "Date") -> pd.DataFrame:
    """Retire toutes les pannes de la semaine ISO indiquée (ex: semaine 24, jugée
    non fiable / à ne pas afficher), partout dans l'application."""
    if date_col not in df.columns or df.empty:
        return df
    iso_week = df[date_col].dt.isocalendar().week
    return df[iso_week != week_number].copy()


def compute_group_totals(df: pd.DataFrame, group_col: str, value_col: str = "Duree_min",
                          agg: str = "sum", top_n: int = None) -> pd.DataFrame:
    """Classement décroissant par groupe, SANS ligne 'Total Result' (celle-ci a été
    retirée de tous les graphiques du dashboard à la demande de l'utilisateur).
    agg='sum' pour la durée totale, agg='count' pour le nombre de pannes."""
    if df.empty or group_col not in df.columns:
        return pd.DataFrame(columns=[group_col, "Valeur"])
    if agg == "count":
        grouped = df.groupby(group_col).size().sort_values(ascending=False)
    else:
        grouped = df.groupby(group_col)[value_col].sum().sort_values(ascending=False)
    result = grouped.reset_index()
    result.columns = [group_col, "Valeur"]
    if top_n:
        result = result.head(top_n)
    return result


def format_occurrence_week(date_value) -> str:
    """Formate une date en 'WKxx-j' (semaine ISO - jour de la semaine, 1=lundi
    ... 7=dimanche), comme dans le plan d'action Excel Versigent
    (ex: WK29-1 = lundi de la semaine 29)."""
    if pd.isna(date_value):
        return ""
    iso = date_value.isocalendar()
    return f"WK{iso.week}-{iso.weekday}"


def format_resp_name(engineer_value) -> str:
    """Reformate 'NOM Prénom 1234' (colonne Principal Engineer) en 'Prénom NOM',
    comme dans la colonne Resp. du plan d'action Excel Versigent
    (ex: 'BOUMDEK Mounir 106' -> 'Mounir BOUMDEK')."""
    if not isinstance(engineer_value, str) or not engineer_value.strip():
        return ""
    tokens = [t for t in engineer_value.strip().split() if not t.isdigit()]
    if len(tokens) >= 2:
        return f"{tokens[1].capitalize()} {tokens[0].upper()}"
    elif tokens:
        return tokens[0]
    return ""


def build_versigent_action_plan(df: pd.DataFrame, machine_col: str = "Machine",
                                 cause_col: str = "Cause", completion_col: str = None,
                                 date_col: str = "Date", engineer_col: str = None,
                                 top_n: int = 8) -> pd.DataFrame:
    """Génère le plan d'action EXACTEMENT dans le format Excel Versigent fourni
    (feuille 'ACTION PLAN') :
      - Item Nr = équipement le plus impactant (top_n par durée cumulée)
      - Problem Description = colonne 'Completion Details' de la panne la plus récente
      - Occurrence Date = 'WKxx-j' (semaine ISO - jour de semaine)
      - Cause = laissée VIDE, à remplir par l'opérateur
      - Action Level = toujours 3
      - Resp = ingénieur responsable (colonne Principal Engineer), reformaté 'Prénom NOM'
      - Status = toujours 'Done'
      - Validation Date = même semaine que Occurrence Date
      - Resp (validation) = toujours 'Ahmed BOUBEGHLI'
      - Status (validation) = toujours 'Done'
      - Effectivity Validation Method = toujours 'Visual Control'
    La colonne 'Action Plan' reste à compléter manuellement. La colonne
    'Date of Implement' a été retirée de ce format à la demande de l'utilisateur."""
    if df.empty:
        return pd.DataFrame(columns=VERSIGENT_ACTION_PLAN_COLUMNS)

    totals = df.groupby(machine_col)["Duree_min"].sum().sort_values(ascending=False)
    top_machines = totals.head(top_n).index.tolist()

    rows = []
    for machine in top_machines:
        sub = df[df[machine_col] == machine].sort_values(date_col, ascending=False)
        if sub.empty:
            continue
        latest = sub.iloc[0]

        problem_desc = ""
        if completion_col and completion_col in sub.columns and pd.notna(latest.get(completion_col)):
            problem_desc = str(latest[completion_col]).strip()

        occ_date = format_occurrence_week(latest[date_col]) if date_col in sub.columns else ""

        resp = ""
        if engineer_col and engineer_col in sub.columns:
            resp = format_resp_name(latest.get(engineer_col))

        rows.append({
            "Item Nr": machine,
            "Problem Description": problem_desc,
            "Occurrence Date": occ_date,
            "Cause": "",  # left blank — filled in manually by the operator
            "Action Plan": "",
            "Action Level": 3,
            "Resp": resp,
            "Status": "Done",
            "Effectivity Validation Method": "Visual Control",
            "Validation Date": occ_date,
            "Resp ": "Ahmed BOUBEGHLI",
            "Status ": "Done",
        })
    return pd.DataFrame(rows, columns=VERSIGENT_ACTION_PLAN_COLUMNS)


def add_repartition_chart_excel_labeled(worksheet, header_row, n_rows, cat_col, val_col, title,
                                         anchor, width=13, height=9, colors=None):
    """Comme add_repartition_chart_excel, mais en PieChart (pas doughnut) avec les
    étiquettes de POURCENTAGE affichées à l'EXTÉRIEUR des parts (position 'outEnd'),
    lisibles même pour les petites parts — demandé explicitement à la place des
    pourcentages illisibles au centre du camembert.

    `colors` (optionnel) : liste de couleurs hex (sans '#') à appliquer part par
    part, dans l'ordre des lignes de données. Permet de reproduire exactement
    les mêmes couleurs que le camembert de l'application (ex: rouge/bleu/vert
    pour le camembert CM/ROB/Ultrasonic) au lieu de la palette orange par défaut."""
    from openpyxl.chart import PieChart, Reference
    from openpyxl.chart.marker import DataPoint
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.chart.label import DataLabelList

    chart = PieChart()
    chart.title = title
    last_row = header_row + n_rows
    val_ref = Reference(worksheet, min_col=val_col, min_row=header_row, max_row=last_row)
    cat_ref = Reference(worksheet, min_col=cat_col, min_row=header_row + 1, max_row=last_row)
    chart.add_data(val_ref, titles_from_data=True)
    chart.set_categories(cat_ref)

    palette = colors or [VERSIGENT_ORANGE, VERSIGENT_BLACK, "E8A867", "595959", "B5651D", "BFBFBF", "8C5A2B", "404040"]
    chart.series[0].data_points = [
        DataPoint(idx=i, spPr=GraphicalProperties(solidFill=palette[i % len(palette)]))
        for i in range(n_rows)
    ]
    chart.dataLabels = DataLabelList()
    chart.dataLabels.showLegendKey = False
    chart.dataLabels.showVal = False
    chart.dataLabels.showCatName = False
    chart.dataLabels.showSerName = False
    chart.dataLabels.showPercent = True
    chart.dataLabels.showBubbleSize = False
    chart.dataLabels.position = "outEnd"

    chart.width = width
    chart.height = height
    worksheet.add_chart(chart, anchor)


def add_bar_chart_excel(worksheet, header_row, n_rows, cat_col, val_col, title, anchor,
                         color=None, width=13, height=8, value_axis_title="Duration (min)"):
    """Graphique à barres simple (sans courbe de cumul), avec l'étiquette de valeur
    affichée AU-DESSUS de chaque barre — utilisé pour les mini-graphiques par
    équipement (CM, ROB, Ultrasonic, Autres) et le graphique Nombre de visites,
    exactement comme sur le dashboard (pas de courbe Pareto sur ces petits
    graphiques, juste les valeurs bien visibles)."""
    from openpyxl.chart import BarChart, Reference
    from openpyxl.chart.label import DataLabelList

    color = color or VERSIGENT_ORANGE
    last_row = header_row + n_rows

    bar = BarChart()
    bar.type = "col"
    bar.title = title
    bar.y_axis.title = value_axis_title
    bar.y_axis.majorGridlines = None
    bar.gapWidth = 40
    bar.x_axis.delete = False

    val_ref = Reference(worksheet, min_col=val_col, min_row=header_row, max_row=last_row)
    cat_ref = Reference(worksheet, min_col=cat_col, min_row=header_row + 1, max_row=last_row)
    bar.add_data(val_ref, titles_from_data=True)
    bar.set_categories(cat_ref)
    bar.series[0].graphicalProperties.solidFill = color
    bar.series[0].graphicalProperties.line.noFill = True

    bar.dataLabels = DataLabelList()
    bar.dataLabels.showLegendKey = False
    bar.dataLabels.showVal = True
    bar.dataLabels.showCatName = False
    bar.dataLabels.showSerName = False
    bar.dataLabels.showPercent = False
    bar.dataLabels.showBubbleSize = False
    bar.dataLabels.position = "outEnd"

    bar.width = width
    bar.height = height
    bar.legend = None
    worksheet.add_chart(bar, anchor)


def build_ishikawa_fig(data, group_col="Machine", value_col="Valeur",
                        effect_label="Total Downtime<br>(Top 10)",
                        title="Ishikawa Diagram — Top 10", height=560):
    """Diagramme d'Ishikawa (arête de poisson) construit à partir du Top 10 du
    Pareto de la semaine (Pareto - Top 10) : chaque équipement du Top 10 devient
    une 'arête' branchée sur l'épine centrale, qui pointe vers l'effet (la boîte
    'Total Downtime') à droite — même palette Versigent (orange/noir) que le
    reste du dashboard."""
    import plotly.graph_objects as go

    items = data.copy().reset_index(drop=True)
    n = len(items)

    fig = go.Figure()
    fig.update_xaxes(visible=False, range=[0, 11.4])
    fig.update_yaxes(visible=False, range=[-3.6, 3.6])
    fig.update_layout(
        title=dict(text=title, font=dict(size=20)),
        height=height, plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(t=70, b=20, l=20, r=20), showlegend=False,
    )

    if n == 0:
        fig.add_annotation(x=5.5, y=0, text="No data", showarrow=False, font=dict(size=16, color="#888"))
        return fig

    spine_x_start, spine_x_end, spine_y = 0.4, 9.2, 0

    # Épine centrale + tête de flèche vers l'effet
    fig.add_shape(type="line", x0=spine_x_start, y0=spine_y, x1=spine_x_end, y1=spine_y,
                  line=dict(color=VERSIGENT_BLACK_HEX, width=3))
    fig.add_annotation(x=spine_x_end, y=spine_y, ax=spine_x_end - 0.5, ay=spine_y,
                        xref="x", yref="y", axref="x", ayref="y", showarrow=True,
                        arrowhead=3, arrowsize=1.6, arrowwidth=3, arrowcolor=VERSIGENT_BLACK_HEX)

    # Boîte "effet" (tête du poisson)
    box_x0, box_x1 = spine_x_end + 0.15, 11.2
    fig.add_shape(type="rect", x0=box_x0, y0=-0.7, x1=box_x1, y1=0.7,
                  line=dict(color=VERSIGENT_ORANGE_HEX, width=2),
                  fillcolor=VERSIGENT_ORANGE_HEX)
    fig.add_annotation(x=(box_x0 + box_x1) / 2, y=0, text=f"<b>{effect_label}</b>",
                        showarrow=False, font=dict(color="white", size=14), align="center")

    # Une arête par équipement du Top 10, alternée au-dessus / en dessous de l'épine,
    # dans l'ordre du classement Pareto (rang 1 = le plus proche de la tête).
    import numpy as np
    xs = np.linspace(spine_x_end - 0.5, spine_x_start + 0.9, n)
    bone_len = 2.5
    for i, x in enumerate(xs):
        name = str(items.loc[i, group_col])
        value = items.loc[i, value_col]
        top = (i % 2 == 0)
        y_end = bone_len if top else -bone_len
        color = VERSIGENT_ORANGE_HEX if top else VERSIGENT_BLACK_HEX
        x_end = x - 0.9

        fig.add_shape(type="line", x0=x_end, y0=y_end, x1=x, y1=spine_y,
                      line=dict(color=color, width=2.5))
        fig.add_annotation(
            x=x_end, y=y_end + (0.35 if top else -0.35),
            text=f"<b>#{i + 1} {name}</b><br>{value:,.0f} min",
            showarrow=False, font=dict(color=color, size=11),
            align="center", yanchor="bottom" if top else "top",
        )
        fig.add_shape(type="circle", x0=x - 0.05, y0=-0.05, x1=x + 0.05, y1=0.05,
                      line=dict(color=color, width=1), fillcolor=color)

    return fig


def add_image_to_sheet(worksheet, image_bytes, anchor="B3", width_px=760, height_px=420):
    """Insère une image PNG générique (ex : diagramme d'Ishikawa exporté en PNG)
    dans une feuille Excel, en la redimensionnant à la taille demandée — même
    logique que add_logo_to_sheet mais réutilisable pour n'importe quelle image
    (pas seulement le logo)."""
    if not image_bytes:
        return
    from io import BytesIO
    try:
        from openpyxl.drawing.image import Image as XLImage
        from PIL import Image as PILImage

        pil_img = PILImage.open(BytesIO(image_bytes))
        pil_img = pil_img.resize((width_px, height_px))
        buf = BytesIO()
        pil_img.save(buf, format="PNG")
        buf.seek(0)
        img = XLImage(buf)
        img.width = width_px
        img.height = height_px
        worksheet.add_image(img, anchor)
    except Exception:
        pass


def add_logo_to_sheet(worksheet, logo_bytes, anchor="A1", height_px=40):
    """Insère le logo Versigent (image) dans le coin de la feuille Excel, pour que
    chaque export (Home, Pareto - Top 10, Action Plan, Data, Detailed Graphs)
    soit visiblement une pièce Versigent, comme demandé."""
    if not logo_bytes:
        return
    from io import BytesIO
    try:
        from openpyxl.drawing.image import Image as XLImage
        from PIL import Image as PILImage

        pil_img = PILImage.open(BytesIO(logo_bytes))
        ratio = pil_img.width / pil_img.height if pil_img.height else 1
        width_px = int(height_px * ratio)
        pil_img = pil_img.resize((width_px, height_px))
        buf = BytesIO()
        pil_img.save(buf, format="PNG")
        buf.seek(0)
        img = XLImage(buf)
        img.width = width_px
        img.height = height_px
        worksheet.add_image(img, anchor)
    except Exception:
        pass
